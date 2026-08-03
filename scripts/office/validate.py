# -*- coding: utf-8 -*-
"""Validador basico de estructura para PPT_Trabajo_Final_Robot_3GDL.pptx.

No requiere MS Office instalado (usa solo python-pptx). Revisa:
  - que el archivo abra sin errores;
  - que ninguna forma se salga de los limites del lienzo (overflow
    estructural real, no aproximacion de renderizado de texto);
  - que las diapositivas de contenido tengan notas del orador;
  - que los numeros de pagina (si existen) sean secuenciales.

Uso:
    python scripts/office/validate.py [ruta_al_pptx]

Si no se pasa ruta, usa 04_presentacion/PPT_Trabajo_Final_Robot_3GDL.pptx
relativo a la raiz del repositorio.

Nota: esto NO detecta texto visualmente cortado dentro de un cuadro (para
eso hay que renderizar a PNG via PowerPoint COM y revisar visualmente,
como se hizo durante el desarrollo de este mazo). Este script solo
atrapa errores estructurales objetivos.
"""
import sys
import os
from pptx import Presentation

DEFAULT_REL_PATH = os.path.join(
    '04_presentacion', 'PPT_Trabajo_Final_Robot_3GDL.pptx'
)


def find_default_path():
    here = os.path.dirname(os.path.abspath(__file__))
    root = os.path.abspath(os.path.join(here, '..', '..'))
    return os.path.join(root, DEFAULT_REL_PATH)


def validate(path):
    errors = []
    warnings = []

    try:
        prs = Presentation(path)
    except Exception as e:
        return [f'No se pudo abrir el archivo: {e}'], []

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    n_slides = len(prs.slides)

    page_numbers = []

    for i, slide in enumerate(prs.slides, 1):
        has_text = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_text = True
            # limites del lienzo
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
            except Exception:
                continue
            if left is None or top is None:
                continue
            if left < -1000 or top < -1000:
                warnings.append(f'Slide {i}: forma "{shape.name}" con posicion negativa ({left},{top})')
            if width and height:
                if (left + width) > slide_w + 5000:
                    errors.append(f'Slide {i}: forma "{shape.name}" se sale del ancho del lienzo '
                                   f'(right={left+width}, limite={slide_w})')
                if (top + height) > slide_h + 5000:
                    errors.append(f'Slide {i}: forma "{shape.name}" se sale del alto del lienzo '
                                   f'(bottom={top+height}, limite={slide_h})')
            # pie de pagina (numero de diapositiva)
            if shape.has_text_frame and abs((shape.left or 0) - 11612880) < 5000 \
                    and abs((shape.top or 0) - 6455664) < 5000:
                t = shape.text_frame.text.strip()
                if t.isdigit():
                    page_numbers.append((i, int(t)))

        if not has_text:
            warnings.append(f'Slide {i}: no tiene ningun texto visible (revisar si es intencional)')

        # notas del orador: solo advertir si falta en una diapositiva de contenido
        # (heuristica: diapositivas de "divisor" suelen tener fondo solido y poco texto)
        notes_text = ''
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if not notes_text:
            warnings.append(f'Slide {i}: sin notas del orador')

    # numeros de pagina deben ser estrictamente crecientes y coincidir con el indice de diapositiva
    for i, n in page_numbers:
        if n != i:
            errors.append(f'Slide {i}: numero de pagina mostrado ({n}) no coincide con su posicion real ({i})')

    return errors, warnings, n_slides


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else find_default_path()
    if not os.path.exists(path):
        print(f'ERROR: no existe el archivo: {path}')
        sys.exit(2)

    result = validate(path)
    errors, warnings = result[0], result[1]
    n_slides = result[2] if len(result) > 2 else None

    print(f'Archivo: {path}')
    if n_slides is not None:
        print(f'Diapositivas: {n_slides}')
    print(f'Errores: {len(errors)}')
    for e in errors:
        print(f'  [ERROR] {e}')
    print(f'Advertencias: {len(warnings)}')
    for w in warnings:
        print(f'  [WARN] {w}')

    if errors:
        sys.exit(1)
    print('OK: validacion estructural basica sin errores.')
    sys.exit(0)


if __name__ == '__main__':
    main()
